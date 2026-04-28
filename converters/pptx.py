from pathlib import Path


def convert_pptx(path: Path) -> str:
    """
    Convierte .pptx a Markdown extrayendo el texto de cada diapositiva.
    Cada diapositiva genera una sección ## con su contenido.
    Las tablas dentro de las diapositivas se convierten a tablas Markdown.
    """
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    sections = []

    for slide_num, slide in enumerate(prs.slides, 1):
        blocks = []

        # Ordenar shapes por posición vertical para respetar el orden visual
        shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))

        for shape in shapes:
            # ── Tablas ────────────────────────────────────────────────────────
            if shape.has_table:
                table = shape.table
                rows  = []
                for row in table.rows:
                    cells = [cell.text.strip().replace('\n', ' ').replace('|', '\\|')
                             for cell in row.cells]
                    rows.append(cells)

                if not rows:
                    continue
                ncols  = max(len(r) for r in rows)
                padded = [r + [''] * (ncols - len(r)) for r in rows]
                md_rows = ['| ' + ' | '.join(padded[0]) + ' |',
                           '| ' + ' | '.join('---' for _ in padded[0]) + ' |']
                for row in padded[1:]:
                    md_rows.append('| ' + ' | '.join(row) + ' |')
                blocks.append('\n'.join(md_rows))
                continue

            # ── Texto ─────────────────────────────────────────────────────────
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Detectar si es título por tamaño de fuente o nivel de outline
                is_title = (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
                            and hasattr(shape, 'placeholder_format')
                            and shape.placeholder_format is not None
                            and shape.placeholder_format.idx == 0)

                level = para.level or 0
                if is_title:
                    blocks.append(f"### {text}")
                elif level == 0 and any(
                    run.font.bold for run in para.runs if run.font.bold is not None
                ):
                    blocks.append(f"**{text}**")
                else:
                    indent = "  " * level
                    blocks.append(f"{indent}- {text}" if level > 0 else text)

        if blocks:
            sections.append(f"## Diapositiva {slide_num}\n\n" + "\n\n".join(blocks))

    return "\n\n---\n\n".join(sections) if sections else ""
